# encoding: utf-8
import io
from pptx import Presentation, util
from PIL import Image


def generate_pptx(rows, columns, images) -> io.BytesIO:
    prs = Presentation()
    w = util.Cm(6)
    h = util.Cm(6)
    left = util.Cm(0.1)
    top = util.Cm(0.1)
    right = util.Cm(0.1)
    bottom = util.Cm(0.1)
    padding_top = util.Cm(2)
    padding_bottom = util.Cm(2)

    prs.slide_width = (w + left + right) * columns
    prs.slide_height = (h + top + bottom) * rows + padding_top + padding_bottom
    blank_slide_layout = prs.slide_layouts[6]

    index = 0
    while index < len(images):
        slide = prs.slides.add_slide(blank_slide_layout)
        for i in range(rows * columns):
            if index >= len(images):
                break
            x = i % columns
            y = i // columns

            with io.BytesIO() as converted_image:
                image_bytes = images[index].stream.read()
                image = Image.open(io.BytesIO(image_bytes))
                print('image format:', image.format)
                # 检查图片格式
                if image.format not in ['BMP', 'GIF', 'JPEG', 'PNG', 'TIFF', 'WMF']:
                    # 转换为 PNG 格式
                    image.save(converted_image, format='JPEG')
                else:
                    converted_image.write(image_bytes)
                converted_image.seek(0)
                # Calculate the aspect ratio
                aspect_ratio = image.width / image.height
                if aspect_ratio > 1:
                    # Landscape orientation
                    new_width = w
                    new_height = w / aspect_ratio
                    x_offset = 0
                    y_offset = (h - new_height) / 2
                else:
                    # Portrait orientation
                    new_height = h
                    new_width = h * aspect_ratio
                    x_offset = (w - new_width) / 2
                    y_offset = 0

                slide.shapes.add_picture(converted_image,
                                         (left + w + right) * x + left + x_offset,
                                         padding_top + (top + h + bottom) * y + top + y_offset,
                                         new_width, new_height)
                index += 1
    ret_bytes = io.BytesIO()
    prs.save(ret_bytes)
    ret_bytes.seek(0)
    return ret_bytes
